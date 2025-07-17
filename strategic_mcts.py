import os
import math
from langchain_core.messages import HumanMessage
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np

# Node for MCTS
class TreeNode:
    def __init__(self, state, parent=None, action=None, depth=0):
        self.state = state
        self.parent = parent
        self.action = action
        self.children = []
        self.visits = 0
        self.reward = 0.0
        self.depth = depth

    def is_leaf(self):
        return len(self.children) == 0

    def add_child(self, child):
        self.children.append(child)

    def uct_score(self, total_simulations, exploration_constant=1.41):
        if self.visits == 0:
            return float("inf")
        exploitation = self.reward / self.visits
        exploration = exploration_constant * math.sqrt(math.log(total_simulations) / self.visits)
        return exploitation + exploration
    
    
def select_node(root):
    node = root
    path = []
    while not node.is_leaf():
        total_visits = sum(child.visits for child in node.children)
        node = max(node.children, key=lambda n: n.uct_score(total_visits))
        path.append(node)
    return node, path

def backpropagate(path, reward):
    for node in path:
        node.visits += 1
        node.reward += reward


from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI

# embedding_model = AzureOpenAIEmbeddings(
#     azure_endpoint="https://retailgenai.openai.azure.com/",
#     api_key="key",
#     api_version="2023-05-15",
#     deployment="text-embedding-ada-002"
# )

# llm = AzureChatOpenAI(
#     api_key="keys",
#     azure_endpoint="https://usedoai4xnaoa01.openai.azure.com/",
#     api_version="2024-12-01-preview",
#     deployment_name="gpt-4o-mini"
# )

from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document

# load RFP vector store
def load_rfp_vectorstore(folder_path):
    vs = FAISS.load_local(folder_path, embedding_model, allow_dangerous_deserialization=True)
    return vs

# load portfolio vector store
def load_portfolio_vectorstore(folder_path):
    vs = FAISS.load_local(folder_path, embedding_model, allow_dangerous_deserialization=True)
    return vs

def strategic_action_retrieve_and_reason(state, llm, rfp_vs, portfolio_vs, top_k=5):
    """
    Combines RFP and portfolio evidence to reason about strategic bid fit.
    """

    # generate search queries
    prompt_q = f"""
        You are a strategy analyst.
        Given the question below, create 3 search queries to retrieve relevant information about:
        1) the RFP requirements
        2) our portfolio and capabilities
        3) relevant lessons from past bids.
        
        Question:
        {state}
    """
    try:
        q_resp = llm.invoke([HumanMessage(content=prompt_q)])
        queries = [q.strip("- ").strip() for q in q_resp.content.strip().split("\n") if q.strip()]
    except Exception as e:
        print(f"❌ could not generate queries: {e}")
        return None

    # retrieve from RFP
    rfp_context = []
    for q in queries:
        try:
            hits = rfp_vs.similarity_search(q, k=top_k)
            rfp_context.extend([doc.page_content for doc in hits])
        except Exception as e:
            print(f"⚠️ rfp retrieval failed: {e}")

    # retrieve from portfolio
    portfolio_context = []
    if portfolio_vs:
        for q in queries:
            try:
                hits = portfolio_vs.similarity_search(q, k=top_k)
                portfolio_context.extend([doc.page_content for doc in hits])
            except Exception as e:
                print(f"⚠️ portfolio retrieval failed: {e}")

    combined_context = "\n\n".join(rfp_context + portfolio_context)

    # now reason
    final_prompt = f"""
        You are a strategic consultant helping to answer the user question below.

        RULES:
        - Base your answer ONLY on the retrieved context below.
        - Avoid hallucinating details not present in the context.
        - If the user requests bullet points, provide them clearly.
        - If the user requests word limits, respect them.
        - Be precise, structured, and professional in tone.

        ====================
        CONTEXT:
        {combined_context}
        ====================

        QUESTION:
        {state}

        ANSWER:
    """

    try:
        answer = llm.invoke([HumanMessage(content=final_prompt)])
        return answer.content.strip()
    except Exception as e:
        print(f"❌ final reasoning failed: {e}")
        return None

def run_strategic_mcts(question, llm, rfp_vs, portfolio_vs, max_iterations=5):
    root = TreeNode(state=question)
    available_actions = [lambda s, l, rv, pv: strategic_action_retrieve_and_reason(s, l, rv, pv)]

    best_answer = None
    best_score = -1

    for _ in range(max_iterations):
        selected, path = select_node(root)
        path = [root] + path

        # expand
        children = []
        for action in available_actions:
            new_state = action(selected.state, llm, rfp_vs, portfolio_vs)
            if new_state:
                child = TreeNode(state=new_state, parent=selected, action=action.__name__, depth=selected.depth+1)
                selected.add_child(child)
                children.append(child)

        if not children:
            continue

        # simulate
        for child in children:
            score = min(len(child.state), 500) / 500.0
            backpropagate(path + [child], score)

            if score > best_score:
                best_score = score
                best_answer = child.state

    return best_answer



SOURCE_DIR = r"C:\Users\YE898BQ\OneDrive - EY\Desktop\strategic-bid-evaluation\MAy 7 RFPs"

def interactive_question_loop():
    # list all folders
    folders = [f for f in os.listdir(SOURCE_DIR) if os.path.isdir(os.path.join(SOURCE_DIR, f))]
    
    print("\nAvailable Folders:\n")
    for idx, folder in enumerate(folders, 1):
        print(f"{idx}. {folder}")
    
    choice = input("\nType the folder number you want to process: ").strip()
    try:
        folder_idx = int(choice) - 1
        folder_name = folders[folder_idx]
    except Exception:
        print("Invalid choice.")
        return
    
    folder_path = os.path.join(SOURCE_DIR, folder_name)
    print(f"\n✅ You chose: {folder_name} ({folder_path})\n")
    
    # chunk and embed once
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    import pdfplumber, docx, openpyxl
    from pptx import Presentation
    from tabulate import tabulate
    
    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=200)
    docs = []
    
    for file in os.listdir(folder_path):
        full_path = os.path.join(folder_path, file)
        ext = file.lower().split(".")[-1]
        title = os.path.basename(file)
        try:
            if ext == "pdf":
                with pdfplumber.open(full_path) as pdf:
                    for page in pdf.pages:
                        raw_text = page.extract_text() or ""
                        for chunk in splitter.split_text(raw_text.strip()):
                            docs.append(Document(page_content=chunk, metadata={"source":title}))
            elif ext == "docx":
                doc = docx.Document(full_path)
                for para in doc.paragraphs:
                    if para.text.strip():
                        for chunk in splitter.split_text(para.text.strip()):
                            docs.append(Document(page_content=chunk, metadata={"source":title}))
            elif ext == "pptx":
                prs = Presentation(full_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            for chunk in splitter.split_text(shape.text.strip()):
                                docs.append(Document(page_content=chunk, metadata={"source":title}))
            elif ext == "xlsx":
                wb = openpyxl.load_workbook(full_path, data_only=True)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        line = " | ".join(str(cell) for cell in row if cell)
                        for chunk in splitter.split_text(line.strip()):
                            docs.append(Document(page_content=chunk, metadata={"source":title}))
        except Exception as e:
            print(f"⚠️ Error: {file} - {e}")
    
    # build vector store once
    print("\n🔄 Building vector store for this folder...")
    vs = FAISS.from_documents(docs, embedding_model)
    vs.save_local(f"./faiss_rfp/{folder_name.replace(' ','_')}")
    print(f"✅ Vector store saved in ./faiss_rfp/{folder_name.replace(' ','_')}")

    # portfolio optional
    portfolio_vs = None

    # interactive question loop
    while True:
        user_q = input("\nType your strategic question (or type 'exit' to quit):\n").strip()
        if user_q.lower() == "exit":
            break
        
        final_answer = run_strategic_mcts(
            question=user_q,
            llm=llm,
            rfp_vs=vs,
            portfolio_vs=portfolio_vs,
            max_iterations=3
        )

        def semantic_factuality_scorer(answer_text, vs, k=5, threshold=0.7):
            """
            Scores how well each reasoning line is supported by
            semantically similar retrieved chunks.
            """
            lines = [line.strip() for line in answer_text.split("\n") if line.strip()]
            supported = 0

            for line in lines:
                try:
                    # embed the line
                    line_vector = embedding_model.embed_query(line)

                    # retrieve chunks
                    hits = vs.similarity_search(line, k=k)

                    # embed the retrieved chunks
                    chunk_vectors = [
                        embedding_model.embed_query(doc.page_content)
                        for doc in hits
                    ]

                    # cosine similarity
                    sims = cosine_similarity(
                        [line_vector],
                        chunk_vectors
                    )[0]

                    # if any match over threshold, count as supported
                    if np.max(sims) >= threshold:
                        supported += 1
                except Exception as e:
                    print(f"⚠️ Semantic scorer failed for line: {line[:40]} — {e}")
                    continue
                
            score = supported / len(lines) if lines else 0
            return round(score, 2)        
        
        credit_score = semantic_factuality_scorer(final_answer, vs)


        print("\n🎯 FINAL STRATEGIC DECISION:\n")
        print(final_answer)
        print(f"\n✅ Credit Score (factuality alignment): {credit_score}\n")


if __name__ == "__main__":
    interactive_question_loop()


